# -*- coding: utf-8 -*-
"""2차 발표 deck 생성 — 이전 자료(IPI_Head_Separation.pptx)의 디자인 시스템을 그대로 사용."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

OUT = (r"C:\Users\Won\Desktop\대학교\AI Secure Lab\내부과제\atlas_poc"
       r"\docs\presentation\IPI_Head_Separation_PoC_2nd [26-08-19].pptx")

# ---------------------------------------------------------------- design tokens
BG        = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x16, 0x18, 0x1D)   # 제목
BODY      = RGBColor(0x56, 0x5C, 0x66)   # 본문
MUTED     = RGBColor(0x8B, 0x90, 0x99)   # 캡션
ORANGE    = RGBColor(0xC0, 0x6A, 0x1F)   # kicker / 강조
BLUE      = RGBColor(0x2F, 0x6F, 0xCE)   # 소제목
RED       = RGBColor(0xC1, 0x44, 0x2A)   # 공격 / 경고
CARD      = RGBColor(0xF2, 0xF2, 0xF4)
CARD_HL   = RGBColor(0xDB, 0xE9, 0xF8)   # 강조 카드
TH_FILL   = RGBColor(0xE9, 0xEA, 0xED)   # 표 헤더
ROW_ALT   = RGBColor(0xF7, 0xF7, 0xF9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LINE      = RGBColor(0xE3, 0xE5, 0xEA)

KR = "Malgun Gothic"
MONO = "Consolas"

M_L, M_W = 0.70, 11.93           # 좌측 마진 / 콘텐츠 폭
Y_KICK, Y_TITLE, Y_BODY = 0.45, 0.85, 1.85
Y_FOOT = 7.05

prs = Presentation()
prs.slide_width = Emu(12191695)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def set_font(run, size, bold=False, color=BODY, font=KR):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    # 한글이 latin typeface만으로는 안 잡혀서 ea/cs도 같이 지정
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def para(tf, spec, size=11.5, color=BODY, bold=False, font=KR,
         first=False, space_before=0, space_after=0, align=PP_ALIGN.LEFT,
         line_spacing=1.25):
    """spec: str 또는 [(text, {size,color,bold,font}), ...]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    chunks = [(spec, {})] if isinstance(spec, str) else spec
    for text, opt in chunks:
        r = p.add_run()
        r.text = text
        set_font(r, opt.get("size", size), opt.get("bold", bold),
                 opt.get("color", color), opt.get("font", font))
    return p


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def card(slide, x, y, w, h, fill=CARD):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.09
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def rect(slide, x, y, w, h, fill, alpha=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    if alpha is not None:
        _alpha(sh, alpha)
    return sh


def oval(slide, x, y, w, h, fill, alpha, line_color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    _alpha(sh, alpha)
    sh.line.color.rgb = line_color
    sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    return sh


def _alpha(shape, pct):
    clr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    el = clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    clr.append(el)


def new_slide(kicker, title, title_size=26):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, BG)
    para(textbox(s, M_L, Y_KICK, M_W, 0.40), kicker,
         size=12, bold=True, color=ORANGE, font=MONO, first=True)
    para(textbox(s, M_L, Y_TITLE, M_W, 0.90), title,
         size=title_size, bold=True, color=INK, first=True, line_spacing=1.1)
    return s


def foot(slide, spec, y=Y_FOOT):
    para(textbox(slide, M_L, y, M_W, 0.40), spec,
         size=10.5, color=MUTED, first=True, line_spacing=1.2)


def table(slide, x, y, w, rows, col_w, row_h=0.36, head_h=0.38,
          sizes=None, aligns=None):
    """rows[0]=헤더. 각 셀은 str 또는 (text, {opt})리스트."""
    n_r, n_c = len(rows), len(col_w)
    h = head_h + row_h * (n_r - 1)
    gf = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    # 기본 스타일 제거(테두리/밴딩) — 셀 채움만 사용
    tbl._tbl.find(qn("a:tblPr")).set("bandRow", "0")

    total = sum(col_w)
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r):
        tbl.rows[i].height = Inches(row_h)

    for ri, row in enumerate(rows):
        for ci, cell_spec in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.045)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = TH_FILL
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else ROW_ALT
            tf = cell.text_frame
            tf.word_wrap = True
            al = (aligns or ["l"] * n_c)[ci]
            alignment = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT,
                         "c": PP_ALIGN.CENTER}[al]
            sz = (sizes or [11] * n_c)[ci]
            if ri == 0:
                para(tf, cell_spec, size=sz, bold=True, color=MUTED,
                     first=True, align=alignment, line_spacing=1.0)
            else:
                para(tf, cell_spec, size=sz, color=BODY,
                     first=True, align=alignment, line_spacing=1.15)
    return gf


def card_head(tf, text, first=False):
    para(tf, text, size=13.5, bold=True, color=BLUE, first=first,
         space_before=0 if first else 9, line_spacing=1.15)


B = lambda t, c=INK: (t, {"bold": True, "color": c})     # noqa: E731
M = lambda t: (t, {"font": MONO, "size": 11})            # noqa: E731


# ================================================================ S1 타이틀
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, BG)
rect(s, 0.70, 2.28, 1.10, 0.055, ORANGE)
para(textbox(s, 0.70, 2.50, 9.5, 0.40), "IPI DEFENSE · FOLLOW-UP REPORT",
     size=13, bold=True, color=ORANGE, font=MONO, first=True)
tf = textbox(s, 0.70, 2.95, 11.5, 2.20)
para(tf, "Read Head, Control Head 분리 PoC", size=34, bold=True,
     color=INK, first=True, line_spacing=1.2)
para(tf, "— 교수님 피드백 대응 결과 (2차)", size=34, bold=True,
     color=INK, line_spacing=1.2)
para(textbox(s, 0.70, 4.55, 10.6, 1.00),
     "AgentDojo 네이티브 평가 도입  ·  head 탐색 방법론 재설계  ·  proxy 지표 재검증",
     size=15, color=BODY, first=True)
para(textbox(s, 0.70, 6.50, 11.0, 0.40),
     "1차 발표 07-28  ·  피드백 07-29  ·  대응 실험 07-29~31  ·  2차 발표 2026-08-19       원종빈",
     size=11.5, color=MUTED, font=MONO, first=True)

# ================================================================ S2 1차 발표 요약
s = new_slide("01 · 이전 내용", "1차 발표 요약 — 가설 · 개입 · 그때의 결과")

card(s, M_L, 1.72, M_W, 0.80)
tf = textbox(s, 1.00, 1.93, 11.3, 0.45)
para(tf, [("[정상 지시] + [외부 데이터 ", {"font": MONO, "size": 13}),
          ("(공격자 명령 삽입)", {"font": MONO, "size": 13, "color": ORANGE}),
          ("] → ", {"font": MONO, "size": 13}),
          ("LLM", {"font": MONO, "size": 13, "bold": True, "color": INK}),
          (" → ", {"font": MONO, "size": 13}),
          ("공격 실행", {"font": MONO, "size": 13, "color": RED})],
     first=True)

card(s, M_L, 2.72, 5.82, 1.32)
tf = textbox(s, 0.95, 2.92, 5.32, 0.95)
card_head(tf, "가설 — Read Head ≠ Control Head", first=True)
para(tf, "외부 데이터의 내용을 읽는 head와, 그 속의 명령을 실행에 옮기는 head가 분리되어 있다",
     size=11.5, space_before=4)

card(s, 6.81, 2.72, 5.82, 1.32)
tf = textbox(s, 7.06, 2.92, 5.32, 0.95)
card_head(tf, "개입 — Edge Knockout", first=True)
para(tf, "head 전체를 끄지 않고, 그 head의 D_inj 방향 attention edge만 pre-softmax에서 차단",
     size=11.5, space_before=4)

table(s, M_L, 4.28, M_W,
      [["데이터", "k=0 공격 성공", "knockout 후", "utility"],
       ["합성 30템플릿 (0.5B~7B)", "0.91 ~ 1.00", [B("0.0000", RED)], "유지 / 소폭 상승"],
       ["합성 unseen 문구 4종", "0.92 ~ 1.00", [B("0.0000", RED)], "유지"],
       ["InjecAgent 1,054개 (1.5B)", "0.1894", "0.0406  (4.7배↓)", "(S9에서 무효 판정)"],
       ["InjecAgent 1,054개 (7B)", "0.5237", "0.1300  (4.0배↓)", "(동일)"]],
      col_w=[3.6, 2.4, 2.7, 3.2], row_h=0.40, aligns=["l", "r", "r", "l"])

foot(s, [("당시 결론: ", {}), B("\"read head와 control head는 분리되어 있다\""),
         ("   —   jaccard(read,·) 0.18~0.38  vs  jaccard(internal,external) 0.48~0.60", {})],
     y=6.42)

# ================================================================ S3 피드백 / 로드맵
s = new_slide("01 · 이전 내용", "받은 피드백 — 세 지시가 사실은 같은 문제의 세 얼굴")

table(s, M_L, Y_BODY, M_W,
      [["피드백 (2026-07-29)", "대응", "슬라이드"],
       [[B("① "), ("AgentDojo로 재현 + utility 측정 제대로", {})],
        "Track A(탐색) / Track B(네이티브 평가) 신규 구축", "S6 ~ S8"],
       [[B("② "), ("키 그룹 2개 vs 데이터셋 모드 4개, read_injected는 왜?", {})],
        "정의 재정리 (문서·코드 주석)", "S4"],
       [[B("③ "), ("attention head 찾는 방법 체계화", {})],
        "자체 진단 5종 + 3소스 비교로 재설계", "S5, S7"]],
      col_w=[5.3, 5.0, 1.6], row_h=0.46, aligns=["l", "l", "c"])

card(s, M_L, 3.85, M_W, 2.35, CARD_HL)
tf = textbox(s, 1.00, 4.08, 11.33, 1.95)
card_head(tf, "①과 ③을 하나로 통합한 이유", first=True)
para(tf, [("1차 발표 P2-d에서 이미 ", {}),
          B("합성 head 단독 / InjecAgent head 단독 / 교집합 세 조건이 거의 동일한 성능"),
          ("을 냈다 (jaccard 0.36).", {})], size=12, space_before=6)
para(tf, "→ \"합성으로 찾은 head를 실제 벤치마크에 전이\"시키는 지금 구조가 합성 데이터의 한계"
         "(노골적 문구·단순 문서 구조)를 그대로 상속할 수 있다는 신호.", size=12, space_before=3)
para(tf, [("→ 그렇다면 AgentDojo는 ", {}), ("단순 재현 대상", {"color": MUTED}),
          ("이 아니라 ", {}), B("head 탐색 소스 자체", ORANGE),
          ("로 삼아야 한다. 두 지시는 같은 작업이었다.", {})], size=12, space_before=3)

foot(s, "아래 S4~S9는 이 로드맵을 ② → ③ → ① 순서로 따라간다.")

# ================================================================ S4 피드백 ②
s = new_slide("02 · 피드백 ② 대응", "\"정상 vs 침입\"이 아니라 \"침입의 두 실행 경로\"")

table(s, M_L, Y_BODY, 6.45,
      [["구분", "개수", "목록"],
       ["relevance 계산 단위(조사)", "3", "read, internal, external"],
       [[B("교집합(선정)에 쓰는 것")], [B("2")], [B("internal, external")]],
       ["key span(창문) 종류", "2", "data_benign, data_inj"],
       [[B("프롬프트 모드")], [B("4")],
        [B("read_clean, read_injected,"), ("\ninternal, external", {"bold": True, "color": INK})]]],
      col_w=[3.3, 0.9, 4.0], row_h=0.46, head_h=0.38,
      sizes=[10.5, 10.5, 10.5], aligns=["l", "c", "l"])

card(s, 7.40, Y_BODY, 5.23, 1.30)
tf = textbox(s, 7.66, 2.03, 4.71, 1.00)
card_head(tf, "internal / external", first=True)
para(tf, [B("둘 다 data_inj를 본다."),
          (" 차이는 backward 타깃뿐 — 자유 텍스트로 새는 경로 vs tool-call로 새는 경로.", {})],
     size=11, space_before=4)

card(s, 7.40, 3.30, 5.23, 1.55)
tf = textbox(s, 7.66, 3.50, 4.71, 1.25)
card_head(tf, "read_injected는 왜 있나", first=True)
para(tf, [("head 탐색에는 ", {}), B("안 쓴다.", RED),
          (" knockout 후 utility 검증 전용. read_clean은 숨은 명령이 없어 끊을 edge 자체가 없어 이 검증에 못 쓴다.", {})],
     size=11, space_before=4)

card(s, M_L, 5.10, M_W, 1.35)
tf = textbox(s, 1.00, 5.30, 11.33, 1.00)
card_head(tf, "다만 — 이 2×2 구조는 합성 데이터셋에서만 성립한다", first=True)
para(tf, "InjecAgent · AgentDojo는 tool-calling 단일 포맷이라 internal/external 구분이 구조상 없다. "
         "→ S7에서 교집합의 기준을 채널 축에서 소스 축으로 교체하게 되는 근거.", size=11.5, space_before=4)

foot(s, "혼동의 원인: \"키 그룹\"이라는 말이 조사 개수(3) · 선정에 쓰는 그룹 수(2) · span 종류(2) · "
        "프롬프트 모드 수(4)를 동시에 가리키고 있었다.")

# ================================================================ S5 피드백 ③ 진단
s = new_slide("02 · 피드백 ③ 대응", "먼저 자체 진단 — 지금까지의 결론이 유효한가")

table(s, M_L, Y_BODY, M_W,
      [["진단 항목", "결과", "판정"],
       ["랜덤 head 기준선",
        [("선정 head(9~14개): k=10 안에 0.9118 → ", {}), B("0.0000", RED),
         ("   /   랜덤 동수: k=40까지도 ", {}), B("0.6448")],
        [("확인", {"bold": True, "color": BLUE})]],
       ["jaccard 우연 기준선",
        [("관측 0.538  vs  우연 기대 0.031  →  ", {}), B("17.6배", ORANGE)],
        [("확인", {"bold": True, "color": BLUE})]],
       ["top-K sweep",
        "K=10(교집합 9개)에서 이미 완전 억제. K=20은 임계점이 아니고 K≥10이면 충분",
        [("확인", {"bold": True, "color": BLUE})]],
       ["dual-use head",
        [("control_heads_both 14개 중 ", {}), B("9개가 read top-20과도 겹침", RED)],
        [("주의", {"bold": True, "color": RED})]],
       ["layer 0 지배",
        [("14개 중 ", {}), B("6개가 layer 0", RED),
         (" (layer 0은 head가 12개뿐 → 절반이 차단됨)", {})],
        [("주의", {"bold": True, "color": RED})]]],
      col_w=[2.6, 8.1, 1.2], row_h=0.44, aligns=["l", "l", "c"])

card(s, M_L, 4.80, M_W, 1.75, CARD_HL)
tf = textbox(s, 1.00, 5.02, 11.33, 1.35)
card_head(tf, "1차 발표의 주장 하나를 철회합니다", first=True)
para(tf, [("\"read head와 control head는 ", {}), ("분리되어 있다", {"color": MUTED}),
          ("\" → ", {}),
          B("\"internal-external끼리의 겹침(0.538)이 read와의 겹침(0.29 / 0.33)보다 뚜렷이 크다\"", INK),
          ("는 상대적 주장까지만 성립.", {})], size=12, space_before=6)
para(tf, "동시에 확인된 것: head 선정 자체는 랜덤 대비 뚜렷이 유효하다 → 합성 데이터셋은 "
         "discovery(탐색) 전용으로 계속 쓰고, 성능 수치는 AgentDojo 네이티브 채점에 맡기기로 결정.",
     size=11.5, space_before=4)

foot(s, "Qwen2.5-1.5B · 템플릿 30개 전체 · 전부 forward-only (비용 거의 없음)")

# ================================================================ S6 Track A/B
s = new_slide("02 · 피드백 ① 대응", "AgentDojo를 어떻게 붙였나 — 탐색과 평가를 분리")

card(s, M_L, Y_BODY, 7.00, 1.85)
tf = textbox(s, 0.92, 2.10, 6.56, 1.40)
para(tf, "                ┌─ Track A (탐색)", size=12, font=MONO, color=INK, first=True,
     line_spacing=1.35)
para(tf, "  AgentDojo ────┤   단일 턴으로 잘라 relevance 계산", size=12, font=MONO, color=BLUE,
     line_spacing=1.35)
para(tf, "                └─ Track B (평가)", size=12, font=MONO, color=INK, line_spacing=1.35)
para(tf, "                    멀티턴 agent loop + 네이티브 채점", size=12, font=MONO, color=ORANGE,
     line_spacing=1.35)

card(s, 8.05, Y_BODY, 4.58, 1.85)
tf = textbox(s, 8.31, 2.05, 4.06, 1.50)
card_head(tf, "왜 나눴나", first=True)
para(tf, "탐색은 backward가 필요해 단일 턴이어야 하고, 평가는 반대로 실제 멀티턴이어야 의미가 있다 "
         "— 두 요구가 양립 불가.", size=11, space_before=3)

card(s, M_L, 3.90, 7.00, 1.55, CARD_HL)
tf = textbox(s, 0.96, 4.10, 6.48, 1.20)
card_head(tf, "네이티브 채점이란", first=True)
para(tf, [("모델 출력 토큰 확률이 아니라, tool을 ", {}),
          B("실제로 실행한 뒤 환경 상태(계좌 잔액·전송된 메일)를 검사하는 결정론적 함수."),
          (" 교수님이 말씀하신 \"제대로 된 성능 측정\"에 대한 답.", {})], size=11.5, space_before=4)

para(textbox(s, 8.05, 3.90, 4.58, 0.30),
     "Track A 어댑터에서 걸린 문제 3가지", size=12, bold=True, color=BLUE, first=True)
table(s, 8.05, 4.28, 4.58,
      [["문제", "해결"],
       ["YAML 개행 정규화", "<INFORMATION> 태그를 앵커로"],
       ["인자 오염형 공격", "해당 case 제외"],
       ["토큰 충돌", "토큰 id 기준 재필터"]],
      col_w=[1.9, 2.7], row_h=0.36, head_h=0.34, sizes=[10, 10])

foot(s, [("Track A: 4개 suite 949조합 → 필터 통과 ", {}), B("220쌍"),
         ("       Track B: KnockoutLocalLLM으로 우리 HF 모델을 agentdojo 파이프라인에 직접 끼워, "
          "롤아웃 내내 D_inj edge를 차단", {})], y=6.30)

# ================================================================ S7 3소스 비교
s = new_slide("02 · head 탐색 재설계", "3소스 head 비교 — 채널 축에서 소스 축으로")

table(s, M_L, Y_BODY, 5.60,
      [["비교", "jaccard", "우연 대비"],
       ["synthetic ↔ InjecAgent", "0.308", [B("12.2배", ORANGE)]],
       ["synthetic ↔ AgentDojo", "0.214", [B("8.5배", ORANGE)]],
       ["InjecAgent ↔ AgentDojo", "0.290", [B("9.5배", ORANGE)]]],
      col_w=[2.9, 1.3, 1.4], row_h=0.40, aligns=["l", "r", "r"])

# --- 3원 벤 다이어그램 (표 아래 / 푸터 위 구간에 배치)
vx, vy, d = 2.00, 3.62, 1.95
oval(s, vx, vy, d, d, BLUE, 22, BLUE)
oval(s, vx + 1.13, vy, d, d, ORANGE, 22, ORANGE)
oval(s, vx + 0.565, vy + 1.00, d, d, RED, 22, RED)
para(textbox(s, 0.72, vy + 0.80, 1.22, 0.30), "synthetic 14",
     size=11, bold=True, color=BLUE, first=True, align=PP_ALIGN.RIGHT)
para(textbox(s, vx + 3.16, vy + 0.80, 1.30, 0.30), "InjecAgent 20",
     size=11, bold=True, color=ORANGE, first=True, align=PP_ALIGN.LEFT)
para(textbox(s, vx + 0.44, vy + 3.03, 2.20, 0.30), "AgentDojo 20",
     size=11, bold=True, color=RED, first=True, align=PP_ALIGN.CENTER)
para(textbox(s, vx + 1.07, vy + 1.13, 0.95, 0.36), "5",
     size=19, bold=True, color=INK, first=True, align=PP_ALIGN.CENTER)

card(s, 6.60, 2.90, 6.03, 1.55)
tf = textbox(s, 6.86, 3.10, 5.51, 1.20)
card_head(tf, "새 질문", first=True)
para(tf, [("교집합의 기준을 채널(internal/external)에서 ", {}), B("소스", ORANGE),
          ("로 바꿨다 — ", {}),
          B("\"어느 소스로 찾은 head가 다른 소스에도 전이되는가.\"", INK)],
     size=11.5, space_before=4)

card(s, 6.60, 4.60, 6.03, 1.85, CARD_HL)
tf = textbox(s, 6.86, 4.82, 5.51, 1.45)
card_head(tf, "3소스 교집합 5개가 전부 layer 0", first=True)
para(tf, [M("(0,1)  (0,3)  (0,6)  (0,7)  (0,10)")], size=11.5, space_before=5)
para(tf, [("S5의 ", {}), B("layer 0 지배", RED),
          (" 우려가 소스를 3개로 늘려도 그대로 유지된다.", {})], size=11.5, space_before=4)

foot(s, "집합 크기 synthetic 14 / InjecAgent 20 / AgentDojo 20  ·  3소스 합집합 36  ·  "
        "세 소스가 무관한 head를 찾은 건 아니지만(우연 대비 8~12배) 같지도 않다.")

# ================================================================ S8 Track B 결과
s = new_slide("02 · Track B 평가", "첫 멀티턴 신호 — 그리고 표본을 늘리려던 두 손잡이")

table(s, M_L, Y_BODY, 6.45,
      [["Qwen2.5-7B (4bit) · 48쌍", "utility", "공격 성공률"],
       ["k=0  (방어 없음)", "0.188", [B("0.021  (1/48)", RED)]],
       ["k=14  (knockout)", "0.188  (변화 없음)", [B("0.000  (0/48)", BLUE)]]],
      col_w=[3.9, 2.2, 2.1], row_h=0.42, aligns=["l", "r", "r"])

card(s, 7.40, Y_BODY, 5.23, 1.22)
tf = textbox(s, 7.66, 2.03, 4.71, 0.95)
card_head(tf, "첫 실제 신호", first=True)
para(tf, "유일하게 성공한 공격(slack/user_task_10)이 knockout으로 실패로 바뀌고, utility는 전혀 안 깎임.",
     size=11, space_before=4)

para(textbox(s, M_L, 3.42, 5.82, 0.30),
     "손잡이 A — 공격 문구가 약해서인가?", size=12.5, bold=True, color=BLUE, first=True)
table(s, M_L, 3.80, 5.82,
      [["공격 기법", "요구하는 것", "성공률"],
       ["important_instructions", "모델이 스스로 추론", "0.021 (1/48)"],
       ["tool_knowledge", "정답 tool·인자를 명시", "0.022 (1/46)"]],
      col_w=[2.5, 1.9, 1.5], row_h=0.38, head_h=0.34,
      sizes=[10, 10, 10], aligns=["l", "l", "r"])
para(textbox(s, M_L, 5.35, 5.82, 0.60),
     [("→ ", {}), B("아니다.", RED), (" 가장 노골적인 문구로 바꿔도 전체 성공률은 그대로.", {})],
     size=11, color=BODY, first=True)

para(textbox(s, 6.81, 3.42, 5.82, 0.30),
     "손잡이 B — 모델이 작아서인가?", size=12.5, bold=True, color=BLUE, first=True)
table(s, 6.81, 3.80, 5.82,
      [["", "7B(4bit)", "14B(4bit)"],
       ["k=0 utility", "0.188", [B("0.267", BLUE)]],
       ["k=0 공격 성공률", "0.021 (1/48)", [B("0.022 (1/45)", RED)]]],
      col_w=[2.4, 1.7, 1.7], row_h=0.38, head_h=0.34,
      sizes=[10, 10, 10], aligns=["l", "r", "r"])
para(textbox(s, 6.81, 5.35, 5.82, 0.60),
     [("→ ", {}), B("절반만 맞다.", RED),
      (" utility는 오르지만(travel 0.00→0.11) 공격 성공률은 불변.", {})],
     size=11, color=BODY, first=True)

foot(s, [B("남은 문제: ", RED),
         ("48쌍 중 성공 공격이 1건이라 통계적으로는 아직 의미 없음. "
          "16GB GPU의 실용 한계(14B-4bit)까지 이미 써버린 상태.", {}),
         ("      4 suite = banking / slack / travel / workspace, suite당 12쌍",
          {"color": MUTED})], y=6.10)

# ================================================================ S9 ★ 핵심 대조
s = new_slide("02 · 핵심 결과", "같은 모델 · 같은 head, 세 가지 측정")

table(s, M_L, 1.78, M_W,
      [["평가 방식", "무엇을 재는가", "k=0 (방어 없음)", "knockout 후"],
       ["합성 30개  (proxy)", "다음 토큰 확률", [B("1.0000", RED)], [B("0.0000", BLUE)]],
       ["InjecAgent 1,054개  (proxy)", "다음 토큰 확률", [B("0.3347", RED)],
        "0.0463  (7.2배↓)"],
       [[B("AgentDojo 45~48쌍  (네이티브)", ORANGE)],
        [B("실제 실행 후 환경 상태", ORANGE)], [B("0.021", RED)], [B("0.000", BLUE)]]],
      col_w=[3.5, 3.3, 2.6, 2.5], row_h=0.44, aligns=["l", "l", "r", "r"])

# --- baseline 공격 성공률 막대 (선형 축)
BAR_X, BAR_MAX, BAR_H = 3.70, 7.20, 0.32
bars = [("합성 (proxy)", 1.0000, RED),
        ("InjecAgent (proxy)", 0.3347, ORANGE),
        ("AgentDojo (네이티브)", 0.0210, BLUE)]
para(textbox(s, M_L, 3.92, 6.0, 0.30),
     "baseline(k=0) 공격 성공률 — 같은 축으로 나란히",
     size=12.5, bold=True, color=INK, first=True)
for i, (label, val, col) in enumerate(bars):
    y = 4.42 + i * 0.62
    para(textbox(s, M_L, y + 0.03, 2.85, 0.30), label,
         size=11, color=BODY, first=True, align=PP_ALIGN.RIGHT)
    rect(s, BAR_X, y, BAR_MAX, BAR_H, RGBColor(0xEC, 0xED, 0xF0))
    rect(s, BAR_X, y, max(BAR_MAX * val, 0.045), BAR_H, col)
    para(textbox(s, BAR_X + BAR_MAX + 0.18, y + 0.03, 1.6, 0.30),
         "%.1f%%" % (val * 100), size=12, bold=True, color=col, first=True)

card(s, M_L, 6.28, M_W, 0.82, CARD_HL)
tf = textbox(s, 1.00, 6.48, 11.33, 0.50)
para(tf, [B("proxy 지표가 실제 공격 성공률을 크게 부풀리고 있었다.", INK),
          ("   차이를 만드는 건 방어 효과가 아니라 ", {}), B("측정 방법", ORANGE),
          ("이다 — \"성능 측정 방법도 제대로 된 걸로\" 지적이 정확했다는 실험적 증거.", {})],
     size=12, first=True)

para(textbox(s, M_L, 7.14, M_W, 0.30),
     "Qwen2.5-14B(bnb-4bit) · 동일한 13개 head · 합성 utility(read)는 0.9360 → 0.9515로 유지",
     size=10.5, color=MUTED, first=True)

# ================================================================ S10 한계
s = new_slide("03 · 이후 진행", "현재 한계 — 다음 단계의 근거")

table(s, M_L, Y_BODY, M_W,
      [["#", "한계", "상태"],
       ["1", [B("통계적 유의성 부족", RED), (" — 48쌍 중 성공 공격 1건", {})], "표본 확대 필요"],
       ["2", [B("layer 0 지배", RED),
              (" — \"명령 인식 회로\" vs \"초기 정보 대역폭 차단\" 미구분", {})], "판별 실험 미설계"],
       ["3", "read / control 완전 분리 아님 (14개 중 9개 dual-use)", "서술 수정 완료"],
       ["4", "합성 데이터셋 content-availability 교란", "discovery 전용으로 봉합"],
       ["5", [("InjecAgent utility 지표는 ", {}), B("인용 불가", RED),
              (" — ASR 감소의 산술적 뒷면", {})], "폐기 결정"],
       ["6", "오라클 스팬 — 배포 시엔 어느 토큰이 주입문인지 모름", "미착수"],
       ["7", "인프라 — 메모리 누수는 우회만, travel suite OOM 잔존", "표본 확대 시 재발 예상"]],
      col_w=[0.6, 8.3, 3.0], row_h=0.44, aligns=["c", "l", "l"])

card(s, M_L, 5.85, M_W, 1.00)
tf = textbox(s, 1.00, 6.05, 11.33, 0.65)
para(tf, [("5번이 ", {}), B("AgentDojo로 가야 했던 핵심 이유"),
          ("입니다 — 피드백 ①의 후반부(\"utility 측정 제대로\")에 대한 저희 쪽 답이기도 합니다.", {})],
     size=12, first=True)

# ================================================================ S11 다음 단계
s = new_slide("03 · 이후 진행", "다음 단계 — 우선순위")

table(s, M_L, Y_BODY, 7.75,
      [["순위", "항목"],
       ["1", [B("표본 확대"), (" (수백 쌍) — travel 배분 축소, 신호 나오는 banking/slack에 재배분", {})]],
       ["2", [B("최종 head 집합 결정"),
              (" — 3소스 교집합 5개(전부 layer 0) vs 각 소스 단독 vs 합집합 36개를 Track B로 비교", {})]],
       ["3", "top-K sweep + random baseline을 3소스 전부에 동일 적용"],
       ["4", [B("layer 0 지배 원인 규명"), (" — 명령 인식인가 정보 대역폭인가", {})]],
       ["5", "AgentDojo baseline 공격 성공률이 왜 2%인지 규명 (모델 크기와 무관함은 확인됨)"],
       ["6", "결과를 methodology.md에 \"Head Selection Methodology\" 절로 통합"]],
      col_w=[0.75, 7.0], row_h=0.52, aligns=["c", "l"], sizes=[10.5, 10.5])

card(s, 8.70, Y_BODY, 3.93, 3.65)
tf = textbox(s, 8.96, 2.07, 3.41, 3.20)
card_head(tf, "중장기", first=True)
for t in ["path patching으로 개별 head 인과성 검증",
          "MMLU 등으로 collateral damage 측정",
          "Llama 계열 교차 검증 (코드는 이미 지원)",
          "오라클 스팬 제거 → 실전 배포 형태",
          "합성 데이터셋 P8 (교란 제거)"]:
    para(tf, [("·  ", {"color": ORANGE}), (t, {})], size=11, space_before=7)

card(s, M_L, 5.72, M_W, 1.12, CARD_HL)
tf = textbox(s, 1.00, 5.94, 11.33, 0.75)
para(tf, [B("2번이 풀리면 4번의 답도 같이 나온다"),
          (" — 3소스 교집합 5개(전부 layer 0)만으로 효과가 나오는지가 갈림길입니다. "
           "효과가 있으면 \"layer 0 대역폭\" 대안 가설이 오히려 강해지고, 없으면 배제됩니다.", {})],
     size=12, first=True)

# ================================================================ S12 마무리
s = new_slide("03 · 마무리", "결론과 논의 요청")

cw, gap = 3.7767, 0.30
for i, (head, lines, col) in enumerate([
    ("한 것", ["피드백 3개 전부 대응",
               "정의 정리 · 자체 진단 5종",
               "AgentDojo Track A·B 신규 구축"], BLUE),
    ("얻은 것", ["첫 멀티턴 네이티브 신호 (1/48 → 0/48)",
                 "3소스 head 비교 (교집합 5개)",
                 "proxy 지표가 부풀려졌다는 정량적 증거"], ORANGE),
    ("남은 것", ["표본 부족 (성공 공격 1건)",
                 "layer 0 지배 미규명",
                 "최종 head 집합 미결정"], RED)]):
    x = M_L + i * (cw + gap)
    card(s, x, Y_BODY, cw, 1.85)
    tf = textbox(s, x + 0.26, 2.05, cw - 0.52, 1.50)
    para(tf, head, size=13.5, bold=True, color=col, first=True)
    for t in lines:
        para(tf, [("·  ", {"color": col}), (t, {})], size=11, space_before=6)

card(s, M_L, 4.00, M_W, 2.42, CARD_HL)
tf = textbox(s, 1.00, 4.22, 11.33, 2.00)
card_head(tf, "여쭙고 싶은 것", first=True)
para(tf, [B("① AgentDojo baseline 공격 성공률 2%를 어떻게 볼 것인가", INK),
          ("  —  (a) 표본 확대에 예산을 더 쓴다   (b) \"이 스케일에선 원래 잘 안 통한다\"를 결과로 보고 "
           "다른 환경을 찾는다  →  ", {}), B("저희 권장은 (a)+(b) 병행", ORANGE)],
     size=11.5, space_before=7)
para(tf, [B("② layer 0 지배를 가장 값싸게 판별할 실험은 무엇일까요", INK)],
     size=11.5, space_before=6)
para(tf, [B("③ \"완전 분리\"를 포기할 때 이 연구의 기여를 무엇으로 서술할까", INK),
          ("  —  현재 후보: 분리가 아니라 ", {}), B("\"편중 + 최소 개입\"", ORANGE)],
     size=11.5, space_before=6)

foot(s, [("이번 사이클의 실질적 성과는 방어 효과를 키운 게 아니라, ", {}),
         B("믿을 수 있는 자(尺)를 갖게 된 것", INK), ("입니다.", {})], y=6.62)

# ----------------------------------------------------------------
prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
