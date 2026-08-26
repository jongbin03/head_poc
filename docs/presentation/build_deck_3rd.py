# -*- coding: utf-8 -*-
"""3차 발표 deck 생성 — 모델 스케일업(Qwen2.5-32B) + 패밀리 교차검증(Llama-3.1-8B).

내용은 IPI_Head_PoC_3rd_script.md를 그대로 옮긴 것. 디자인 토큰/헬퍼는
build_deck_2nd.py와 동일 — import하면 2차 deck이 재빌드되므로(그쪽도 스크립트) 의도적으로
복제했다 (build_deck_addendum_0819.py와 같은 이유).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

OUT = (r"C:\Users\Won\Desktop\대학교\AI Secure Lab\내부과제\atlas_poc"
       r"\docs\presentation\IPI_Head_Separation_PoC_3rd [26-08-26].pptx")

# ---------------------------------------------------------------- design tokens
BG        = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x16, 0x18, 0x1D)   # 제목
BODY      = RGBColor(0x56, 0x5C, 0x66)   # 본문
MUTED     = RGBColor(0x8B, 0x90, 0x99)   # 캡션
ORANGE    = RGBColor(0xC0, 0x6A, 0x1F)   # kicker / 강조
BLUE      = RGBColor(0x2F, 0x6F, 0xCE)   # 소제목
RED       = RGBColor(0xC1, 0x44, 0x2A)   # 공격 / 경고
CARD      = RGBColor(0xF2, 0xF2, 0xF4)
CARD_HL   = RGBColor(0xDB, 0xE9, 0xF8)   # 강조 카드
TH_FILL   = RGBColor(0xE9, 0xEA, 0xED)   # 표 헤더
ROW_ALT   = RGBColor(0xF7, 0xF7, 0xF9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LINE      = RGBColor(0xE3, 0xE5, 0xEA)

KR = "Malgun Gothic"
MONO = "Consolas"

M_L, M_W = 0.70, 11.93           # 좌측 마진 / 콘텐츠 폭
Y_KICK, Y_TITLE, Y_BODY = 0.45, 0.85, 1.85
Y_FOOT = 7.05

prs = Presentation()
prs.slide_width = Emu(12191695)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def set_font(run, size, bold=False, color=BODY, font=KR):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def para(tf, spec, size=11.5, color=BODY, bold=False, font=KR,
         first=False, space_before=0, space_after=0, align=PP_ALIGN.LEFT,
         line_spacing=1.25):
    """spec: str 또는 [(text, {size,color,bold,font}), ...]"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    chunks = [(spec, {})] if isinstance(spec, str) else spec
    for text, opt in chunks:
        r = p.add_run()
        r.text = text
        set_font(r, opt.get("size", size), opt.get("bold", bold),
                 opt.get("color", color), opt.get("font", font))
    return p


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def card(slide, x, y, w, h, fill=CARD):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.09
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def rect(slide, x, y, w, h, fill, alpha=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    if alpha is not None:
        _alpha(sh, alpha)
    return sh


def _alpha(shape, pct):
    clr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
    el = clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))})
    clr.append(el)


def new_slide(kicker, title, title_size=26):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.33, 7.5, BG)
    para(textbox(s, M_L, Y_KICK, M_W, 0.40), kicker,
         size=12, bold=True, color=ORANGE, font=MONO, first=True)
    para(textbox(s, M_L, Y_TITLE, M_W, 0.90), title,
         size=title_size, bold=True, color=INK, first=True, line_spacing=1.1)
    return s


def foot(slide, spec, y=Y_FOOT):
    para(textbox(slide, M_L, y, M_W, 0.40), spec,
         size=10.5, color=MUTED, first=True, line_spacing=1.2)


def table(slide, x, y, w, rows, col_w, row_h=0.36, head_h=0.38,
          sizes=None, aligns=None):
    """rows[0]=헤더. 각 셀은 str 또는 (text, {opt})리스트."""
    n_r, n_c = len(rows), len(col_w)
    h = head_h + row_h * (n_r - 1)
    gf = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    tbl._tbl.find(qn("a:tblPr")).set("bandRow", "0")

    total = sum(col_w)
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Emu(int(Inches(w) * cw / total))
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r):
        tbl.rows[i].height = Inches(row_h)

    for ri, row in enumerate(rows):
        for ci, cell_spec in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.045)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = TH_FILL
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else ROW_ALT
            tf = cell.text_frame
            tf.word_wrap = True
            al = (aligns or ["l"] * n_c)[ci]
            alignment = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT,
                         "c": PP_ALIGN.CENTER}[al]
            sz = (sizes or [11] * n_c)[ci]
            if ri == 0:
                para(tf, cell_spec, size=sz, bold=True, color=MUTED,
                     first=True, align=alignment, line_spacing=1.0)
            else:
                para(tf, cell_spec, size=sz, color=BODY,
                     first=True, align=alignment, line_spacing=1.15)
    return gf


def card_head(tf, text, first=False):
    para(tf, text, size=13.5, bold=True, color=BLUE, first=first,
         space_before=0 if first else 9, line_spacing=1.15)


B = lambda t, c=INK: (t, {"bold": True, "color": c})     # noqa: E731
M = lambda t: (t, {"font": MONO, "size": 11})            # noqa: E731


# ================================================================ S1 타이틀
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, BG)
rect(s, 0.70, 2.28, 1.10, 0.055, ORANGE)
para(textbox(s, 0.70, 2.50, 9.5, 0.40), "IPI DEFENSE · SCALE & FAMILY VALIDATION",
     size=13, bold=True, color=ORANGE, font=MONO, first=True)
tf = textbox(s, 0.70, 2.95, 11.5, 2.20)
para(tf, "Read Head, Control Head 분리 PoC", size=34, bold=True,
     color=INK, first=True, line_spacing=1.2)
para(tf, "— 확장 검증 (3차)", size=34, bold=True, color=INK, line_spacing=1.2)
para(textbox(s, 0.70, 4.55, 10.6, 1.00),
     "모델 스케일업 (Qwen2.5-32B)  ·  모델 패밀리 교차검증 (Llama-3.1-8B)",
     size=15, color=BODY, first=True)
para(textbox(s, 0.70, 6.50, 11.0, 0.40),
     "2차 발표 08-19  ·  서버 이전 + 확장 실험 08-21~26  ·  3차 발표 2026-08-26       원종빈",
     size=11.5, color=MUTED, font=MONO, first=True)

# ================================================================ S2 서론
s = new_slide("01 · 서론", "이번 사이클에 검증한 것")

card(s, M_L, 1.85, M_W, 1.05)
tf = textbox(s, 0.96, 2.05, 11.3, 0.75)
para(tf, [("2차 발표 결론: knockout(head의 D", {}), ("inj", {"size": 9}),
          (" edge만 차단)이 7B에서 ", {}), B("공격을 억제"), (" 하면서 ", {}),
          B("정상 기능(utility)은 유지", BLUE), ("했다.", {})], size=13, first=True)

para(textbox(s, M_L, 3.20, M_W, 0.50),
     [("이 효과가 ", {}), B("① 더 큰 모델", ORANGE), (" 에서도, ", {}),
      B("② 다른 모델 계열", ORANGE), (" 에서도 재현되는가?", {})],
     size=17, bold=True, color=INK, first=True)

cw, gap = 5.815, 0.30
card(s, M_L, 4.05, cw, 1.85, CARD_HL)
tf = textbox(s, M_L + 0.26, 4.27, cw - 0.52, 1.50)
card_head(tf, "① 스케일 축", first=True)
para(tf, "Qwen2.5-7B → Qwen2.5-32B (4bit)", size=12.5, space_before=6)
para(tf, "같은 모델 계열, 크기만 키워서 같은 실험 반복", size=11.5, color=MUTED, space_before=4)

card(s, M_L + cw + gap, 4.05, cw, 1.85, CARD_HL)
tf = textbox(s, M_L + cw + gap + 0.26, 4.27, cw - 0.52, 1.50)
card_head(tf, "② 패밀리 축", first=True)
para(tf, "Qwen2.5 → Llama-3.1-8B (~7-8B 스케일 고정)", size=12.5, space_before=6)
para(tf, "다른 아키텍처에서 같은 head 탐색·knockout 파이프라인 재실행", size=11.5,
     color=MUTED, space_before=4)

foot(s, "두 축 모두 코드 수정 없이 동일한 파이프라인(head 탐색 → 검증)을 재사용한다.")

# ================================================================ S3 방법 - 환경
s = new_slide("02 · 실험 방법", "실험 환경")

table(s, M_L, Y_BODY, M_W,
      [["항목", "내용"],
       ["서버", "SSH 공용 서버, Titan RTX 24GB ×3"],
       ["dtype", "bf16 (Turing 에뮬레이션, fp32 대비 오차 무시할 수준)"],
       ["스케일 축 모델", [B("Qwen2.5-32B-Instruct"), (" (4bit)", {})]],
       ["패밀리 축 모델", [B("Llama-3.1-8B-Instruct"), (" (fp16)", {})]],
       ["벤치마크", "AgentDojo — banking / slack / travel / workspace 4개 suite"]],
      col_w=[3.2, 8.73], row_h=0.62, head_h=0.42, sizes=[12, 12])

foot(s, "실제 tool 실행 기반 멀티턴 벤치마크(AgentDojo)를 두 축 모두 동일하게 사용.")

# ================================================================ S4 방법 - head 탐색
s = new_slide("02 · 실험 방법", "Head 탐색 (Track A)")

card(s, M_L, Y_BODY, M_W, 1.10)
tf = textbox(s, 0.96, 2.05, 11.3, 0.80)
para(tf, [("relevance", {"font": MONO, "bold": True, "color": INK}),
          (" = attention weight × gradient  —  ", {}),
          B("\"주입된 명령을 실제로 본 head\""),
          (" 를 랭킹, 상위 20개 추출 (AttnLRP 기반)", {})], size=13, first=True)

card(s, M_L, 3.20, 5.815, 2.00, CARD_HL)
tf = textbox(s, M_L + 0.26, 3.42, 5.28, 1.65)
card_head(tf, "user_task 단위 group split", first=True)
para(tf, "같은 과업(user_task)의 다른 injection이 탐색·평가 양쪽에 섞이지 않게 "
         "통째로 한쪽에만 배정", size=11.5, space_before=6)
para(tf, "→ 데이터 누수 방지", size=11.5, color=BLUE, space_before=4)

card(s, M_L + 5.815 + 0.30, 3.20, 5.815, 2.00, CARD_HL)
tf = textbox(s, M_L + 5.815 + 0.30 + 0.26, 3.42, 5.28, 1.65)
card_head(tf, "suite별 층화 샘플링", first=True)
para(tf, "4개 suite가 고르게 표본에 들어가도록 쿼터를 배정 — 특정 suite가 표본을 "
         "독점하지 않게", size=11.5, space_before=6)
para(tf, "→ suite별 편차는 결과에 그대로 공개", size=11.5, color=BLUE, space_before=4)

foot(s, "탐색은 단일 턴(주입 직후 응답)만 잘라서 backward를 계산 — 검증(Track B)과 요구사항이 달라 분리.")

# ================================================================ S5 방법 - 평가
s = new_slide("02 · 실험 방법", "평가 (Track B)")

card(s, M_L, Y_BODY, M_W, 1.30)
tf = textbox(s, 0.96, 2.06, 11.3, 1.00)
para(tf, [("AgentDojo의 ", {}), B("실제 멀티턴 agent loop"),
          (" 안에 모델을 직접 끼워 넣고 실행 — 토큰 확률 proxy가 아니라 ", {}),
          B("tool을 실제로 실행한 뒤 환경 상태로 채점", ORANGE),
          (" (네이티브 채점)", {})], size=13, first=True)

table(s, M_L, 3.55, 7.20,
      [["조건", "설명"],
       ["k = 0", "방어 없음 (baseline)"],
       ["k = N", "찾은 head들의 D_inj 방향 attention edge만 차단 (knockout)"]],
      col_w=[1.6, 5.6], row_h=0.52, head_h=0.38, sizes=[12, 11.5])

card(s, 8.20, 3.55, 4.43, 2.00)
tf = textbox(s, 8.46, 3.75, 3.91, 1.65)
card_head(tf, "두 지표", first=True)
para(tf, [B("utility", BLUE), (" — 원래 과업 성공 여부", {})], size=11.5, space_before=6)
para(tf, [B("security", RED), (" — 공격 성공 여부 (낮을수록 좋음)", {})], size=11.5, space_before=4)

foot(s, "head 탐색에 쓰인 user_task는 평가 후보에서 제외(held-out) — 같은 데이터로 찾고 검증하지 않는다.")

# ================================================================ S6 결과 개요 - 모델 x 공격 비교
s = new_slide("03 · 결과", "결과 개요 — 모델 × 공격 종류 비교")

table(s, M_L, Y_BODY, M_W,
      [["모델", "공격", "n", "utility  k=0", "utility  k=N", "ASR  k=0", "ASR  k=N"],
       ["Qwen2.5-7B", "important_instructions", "48", "18.75%", "18.75%", "2.08%", "0%"],
       ["Qwen2.5-7B", "tool_knowledge", "46", "19.6%", "17.4%", "2.17%", "0%"],
       [[B("Qwen2.5-32B", BLUE)], "important_instructions", "57", "45.6%", "47.4%",
        [B("8.77%", RED)], [B("3.51%", RED)]],
       [[B("Qwen2.5-32B", BLUE)], [B("tool_knowledge", ORANGE)], "55", "54.5%", "56.4%",
        [B("14.5%", RED)], [B("9.1%", RED)]],
       [[B("Llama-3.1-8B", BLUE)], "important_instructions", "44", "38.6%", "38.6%",
        [B("4.5%", RED)], [B("0%", RED)]],
       [[B("Llama-3.1-8B", BLUE)], [B("tool_knowledge", ORANGE)], "44", "36.4%", "34.1%",
        [B("4.5%", RED)], [B("0%", RED)]]],
      col_w=[1.95, 2.65, 0.65, 1.55, 1.55, 1.55, 1.83], row_h=0.42, head_h=0.38,
      sizes=[10, 9.5, 10, 9.5, 9.5, 9.5, 9.5], aligns=["l", "l", "c", "r", "r", "r", "r"])

card(s, M_L, 5.35, M_W, 1.30, CARD_HL)
tf = textbox(s, 1.00, 5.55, 11.33, 1.00)
para(tf, [("세 모델·두 공격 전부 공격 성공률이 knockout 후 ", {}),
          B("방향이 일관되게 낮아진다", BLUE), (" (억제).", {})], size=12, first=True)
para(tf, [B("⚠️ Qwen2.5-7B 두 행은 구버전 파이프라인", RED),
          ("(held-out split 도입 전) 결과라 32B/Llama-8B와 직접 비교는 caveat 필요 — "
           "정량 비교는 다음 슬라이드부터(같은 파이프라인끼리). Llama 두 행은 travel 제외(n=44).", {})],
     size=11, space_before=5)

foot(s, "Llama-3.1-8B 두 행 모두 travel suite 제외(파서 미지원, S9/S10 각주 참고).")

# ================================================================ S7 결과 - 32B important_instructions
s = new_slide("03 · 결과", "32B / important_instructions (n=57)")

table(s, M_L, Y_BODY, M_W,
      [["suite", "n", "utility  k=0", "utility  k=N", "ASR  k=0", "ASR  k=N"],
       ["banking", "15", "53.3%", "73.3%", "0%", "0%"],
       [[B("slack")], "15", "53.3%", "53.3%", [B("33.3%", RED)], [B("13.3%", RED)]],
       ["travel", "15", "20.0%", "6.7%", "0%", "0%"],
       ["workspace", "12", "58.3%", "58.3%", "0%", "0%"],
       [[B("합계")], [B("57")], [B("45.6%", BLUE)], [B("47.4%", BLUE)],
        [B("8.77%", RED)], [B("3.51%", RED)]]],
      col_w=[2.2, 0.8, 2.1, 2.1, 2.2, 2.53], row_h=0.44, aligns=["l", "c", "r", "r", "r", "r"])

card(s, M_L, 5.10, M_W, 1.10, CARD_HL)
tf = textbox(s, 1.00, 5.30, 11.33, 0.75)
para(tf, [("utility는 knockout 후에도 ", {}), B("유지(소폭 상승)", BLUE),
          (", 공격 성공률은 ", {}), B("약 60% 상대 감소", RED),
          (". 공격 성공 사례는 slack에 집중 — 나머지 3개 suite는 baseline부터 0%.", {})],
     size=12, first=True)

foot(s, "results/2026-08-24_s4_32b/agentdojo_eval.json · --eval_split heldout")

# ================================================================ S8 결과 - 32B tool_knowledge
s = new_slide("03 · 결과", "32B / tool_knowledge (n=55)")

table(s, M_L, Y_BODY, M_W,
      [["suite", "n", "utility  k=0", "utility  k=N", "ASR  k=0", "ASR  k=N"],
       ["banking", "15", "86.7%", "80.0%", "0%", "0%"],
       [[B("slack")], "15", "46.7%", "53.3%", [B("53.3%", RED)], [B("33.3%", RED)]],
       ["travel", "12", "25.0%", "25.0%", "0%", "0%"],
       ["workspace", "13", "53.8%", "61.5%", "0%", "0%"],
       [[B("합계")], [B("55")], [B("54.5%", BLUE)], [B("56.4%", BLUE)],
        [B("14.5%", RED)], [B("9.1%", RED)]]],
      col_w=[2.2, 0.8, 2.1, 2.1, 2.2, 2.53], row_h=0.44, aligns=["l", "c", "r", "r", "r", "r"])

card(s, M_L, 5.10, M_W, 1.10, CARD_HL)
tf = textbox(s, 1.00, 5.30, 11.33, 0.75)
para(tf, [("공격 문구를 강화하자 baseline ASR이 ", {}), B("8.77%→14.5%로 뜀", RED),
          (" — 그래도 knockout 후 억제 방향은 유지. utility는 오히려 소폭 상승(54.5%→56.4%).", {})],
     size=12, first=True)

foot(s, "results/2026-08-24_s4_32b/agentdojo_eval_toolknowledge.json · --eval_split heldout")

# ================================================================ S9 결과 - Llama important_instructions
s = new_slide("03 · 결과", "Llama-3.1-8B / important_instructions (n=44)")

table(s, M_L, Y_BODY, M_W,
      [["suite", "n", "utility  k=0", "utility  k=N", "ASR  k=0", "ASR  k=N"],
       ["banking", "15", "40.0%", "40.0%", "0%", "0%"],
       [[B("slack")], "15", "53.3%", "53.3%", [B("13.3%", RED)], [B("0%", RED)]],
       ["workspace", "14", "21.4%", "21.4%", "0%", "0%"],
       [[B("합계")], [B("44")], [B("38.6%", BLUE)], [B("38.6%", BLUE)],
        [B("4.5%", RED)], [B("0%", RED)]]],
      col_w=[2.2, 0.8, 2.1, 2.1, 2.2, 2.53], row_h=0.44, aligns=["l", "c", "r", "r", "r", "r"])

card(s, M_L, 4.75, M_W, 1.10, CARD_HL)
tf = textbox(s, 1.00, 4.95, 11.33, 0.75)
para(tf, [("utility ", {}), B("완전 유지", BLUE), (", 공격 성공률 ", {}),
          B("완전 억제(4.5%→0%)", RED),
          (" — Qwen2가 아닌 다른 아키텍처에서도 같은 패턴이 재현됨.", {})],
     size=12, first=True)

foot(s, "(각주) travel suite는 이번 파이프라인이 멀티 tool-call 응답 파싱을 지원하지 않아 제외 — "
        "실제 성능 결과 아님. results/2026-08-25_s6_llama8b/agentdojo_eval.json", y=6.20)

# ================================================================ S10 결과 - Llama tool_knowledge
s = new_slide("03 · 결과", "Llama-3.1-8B / tool_knowledge (n=44)")

table(s, M_L, Y_BODY, M_W,
      [["suite", "n", "utility  k=0", "utility  k=N", "ASR  k=0", "ASR  k=N"],
       ["banking", "15", "33.3%", "26.7%", "0%", "0%"],
       [[B("slack")], "15", "53.3%", "53.3%", [B("13.3%", RED)], [B("0%", RED)]],
       ["workspace", "14", "21.4%", "21.4%", "0%", "0%"],
       [[B("합계")], [B("44")], [B("36.4%", BLUE)], [B("34.1%", BLUE)],
        [B("4.5%", RED)], [B("0%", RED)]]],
      col_w=[2.2, 0.8, 2.1, 2.1, 2.2, 2.53], row_h=0.44, aligns=["l", "c", "r", "r", "r", "r"])

card(s, M_L, 4.75, M_W, 1.10, CARD_HL)
tf = textbox(s, 1.00, 4.95, 11.33, 0.75)
para(tf, [("공격 성공률은 여기서도 ", {}), B("완전 억제(4.5%→0%)", RED),
          (" — important_instructions와 동일. utility는 소폭 하락(36.4%→34.1%, banking "
           "1건이 knockout 후 실패로 전환) — 32B만큼 완전하진 않지만 큰 폭 유지.", {})],
     size=12, first=True)

foot(s, "(각주) travel suite 제외 — 위와 동일 사유. "
        "results/2026-08-25_s6_llama8b/agentdojo_eval_toolknowledge.json", y=6.20)

# ================================================================ S11 결과 - 공격 강도 비교(종합)
s = new_slide("03 · 결과", "공격 강도 비교 — 모델 간 종합")

table(s, M_L, Y_BODY, M_W,
      [["모델", "공격", "n", "ASR  k=0", "ASR  k=N", "상대 감소율"],
       [[B("Qwen2.5-32B", BLUE)], "important_instructions", "57", "8.77%", "3.51%", "60%"],
       [[B("Qwen2.5-32B", BLUE)], [B("tool_knowledge", ORANGE)], "55", "14.5%", "9.1%", "37.5%"],
       [[B("Llama-3.1-8B", BLUE)], "important_instructions", "44", "4.5%", "0%",
        [B("100%", BLUE)]],
       [[B("Llama-3.1-8B", BLUE)], [B("tool_knowledge", ORANGE)], "44", "4.5%", "0%",
        [B("100%", BLUE)]]],
      col_w=[2.3, 3.1, 0.7, 1.9, 1.9, 2.03], row_h=0.46, head_h=0.40,
      sizes=[11, 10, 11, 10.5, 10.5, 11], aligns=["l", "l", "c", "r", "r", "r"])

card(s, M_L, 3.95, M_W, 1.85, CARD_HL)
tf = textbox(s, 1.00, 4.17, 11.33, 1.50)
para(tf, [B("32B: ", BLUE), ("공격이 강해질수록 knockout 후에도 남는 절대적 위험이 커짐"
          "(3.51%→9.1%) — 억제는 하지만 완전하진 않음.", {})], size=12, first=True)
para(tf, [B("Llama-3.1-8B: ", BLUE),
          ("공격 강도와 무관하게 두 공격 다 ", {}), B("완전 억제(100%)", RED),
          (" — 표본이 작아(성공 사례 2건뿐) 과대해석은 주의해야 하지만, 32B와 다른 양상이라는 "
           "점은 흥미로운 대조.", {})], size=12, space_before=8)

foot(s, "results/2026-08-24_s4_32b/*.json · results/2026-08-25_s6_llama8b/*.json")

# ================================================================ S12 한계 및 다음 실험
s = new_slide("04 · 한계", "한계 및 다음 실험")

table(s, M_L, Y_BODY, M_W,
      [["#", "한계", "향후 방향"],
       ["1", [B("표본 크기", RED),
              (" — suite마다 무작위 15쌍만 평가, 전수 평가 아님 (예: workspace는 실제 "
               "후보 pool 300~450쌍 중 15개만 사용)", {})],
        "suite별 전체 평가로 확대"],
       ["2", [B("suite별 편차", RED),
              (" — 공격 성공 사례가 거의 slack에만 몰려 있음(나머지 3개 suite는 baseline"
               "부터 0%인 경우가 많음)", {})],
        "결과를 \"slack 주도\"로 caveat, 더 다양한 공격 시나리오 필요"],
       ["3", [B("AttnLRP(lxt) 검증 아키텍처 범위", RED),
              (" — head 탐색(backward LRP)이 정상 동작을 확인한 건 지금까지 "
               "Qwen2.5·Llama-3.1 두 계열뿐", {})],
        [("같은 표준 구조(RMSNorm+SwiGLU+표준 attention)를 쓰는 다른 모델(Mistral, "
          "DeepSeek dense 계열 등)로 ", {}), B("확장이 코드상 저비용으로 가능해 보임", BLUE),
         (" — 다음 실험 후보로 제안", {})]],
       ["4", [B("tool-call 파싱", RED),
              (" — 100% 성공은 아님(모델별 70%대) — 실패한 턴은 tool 미실행으로 집계돼 "
               "utility가 과소평가될 수 있음", {})],
        "모델별 파서 정교화 지속"]],
      col_w=[0.5, 7.6, 3.83], row_h=0.82, aligns=["c", "l", "l"], sizes=[11, 11, 10.5])

# ================================================================ S13 마무리
s = new_slide("05 · 마무리", "결론")

cw, gap = 3.7767, 0.30
for i, (head, lines, col) in enumerate([
    ("스케일 축", ["Qwen2.5 7B → 32B",
                  "utility 유지 + 공격 60% 감소"], BLUE),
    ("패밀리 축", ["Qwen2 → Llama-3.1",
                  "utility 유지 + 공격 완전 억제"], BLUE),
    ("공격 강도 축", ["문구를 더 명확하게",
                    "32B는 잔여 위험 증가, Llama는 완전 억제 유지"], ORANGE)]):
    x = M_L + i * (cw + gap)
    card(s, x, Y_BODY, cw, 1.65)
    tf = textbox(s, x + 0.26, 2.05, cw - 0.52, 1.30)
    para(tf, head, size=13.5, bold=True, color=col, first=True)
    for t in lines:
        para(tf, [("·  ", {"color": col}), (t, {})], size=11.5, space_before=8)

card(s, M_L, 3.85, M_W, 1.15, CARD_HL)
tf = textbox(s, 1.00, 4.07, 11.33, 0.80)
para(tf, [B("세 축 모두 \"공격 억제 + 정상 기능 보존\" 패턴이 재현됨", INK),
          (" — 7B 단일 모델·단일 공격 문구에서만 성립하던 효과가 아니었다.", {})],
     size=13, first=True)

foot(s, [B("한계: ", RED),
         ("표본 크기 · suite 편차 · lxt 검증 아키텍처 범위 — 상세는 이전 슬라이드", {})],
     y=5.55)

# ----------------------------------------------------------------
prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
